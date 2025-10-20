from __future__ import annotations
"""
Create a client-facing PowerPoint deck (non-technical) for RAMA Agentic AI.
Generates ~15 slides with a clean, modern look (green/earth palette).

Output: client_deck_rama.pptx in project root.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# Palette (campo / natural)
GREEN_DARK = RGBColor(0x3d, 0x74, 0x35)
GREEN_LIGHT = RGBColor(0x8f, 0xcb, 0x7f)
EARTH = RGBColor(0xc5, 0xac, 0x85)
BG_LIGHT = RGBColor(0xf7, 0xfd, 0xf5)


def _add_banner(slide, title: str, subtitle: str | None = None):
    """Add a top banner with brand color and centered title/subtitle."""
    shapes = slide.shapes
    # background panel (light) using slide background (no private attrs)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT

    # top banner
    banner = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.2))
    banner.fill.solid()
    banner.fill.fore_color.rgb = GREEN_DARK
    banner.line.fill.background()

    # title text box
    tx = shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12.2), Inches(0.9))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle:
        sub = shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.2), Inches(0.5))
        sub_tf = sub.text_frame
        sub_tf.clear()
        ps = sub_tf.paragraphs[0]
        ps.text = subtitle
        ps.font.size = Pt(18)
        ps.font.color.rgb = GREEN_DARK


def _add_bullets(slide, items: list[str], y_start_inches: float = 1.9, font_size: int = 20):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(y_start_inches), Inches(12.0), Inches(5.8))
    tf = box.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = it
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def _add_block_text(slide, text: str, y_start_inches: float = 1.9, font_size: int = 18, mono: bool = False):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(y_start_inches), Inches(12.0), Inches(6.0))
    tf = box.text_frame
    tf.clear()
    # Preserve line breaks
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(font_size)
        if mono:
            p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def add_slide_title_bullets(prs: Presentation, title: str, bullets: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    _add_banner(slide, title, subtitle)
    _add_bullets(slide, bullets)
    return slide


def add_slide_title_text(prs: Presentation, title: str, text: str, subtitle: str | None = None, mono: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_banner(slide, title, subtitle)
    _add_block_text(slide, text, mono=mono)
    return slide


def build_client_deck() -> Presentation:
    prs = Presentation()
    # Widescreen 16:9 is default; keep it

    # 1) Cover
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_banner(slide, "RAMA Agentic AI", "Gestión integral de propiedades con IA")
    # Brand subtitle footer
    foot = slide.shapes.add_textbox(Inches(0.9), Inches(6.8), Inches(10.0), Inches(0.6))
    f_tf = foot.text_frame
    f_tf.text = "Documentos · Números · Resumen · Conversación natural"
    f_tf.paragraphs[0].font.size = Pt(18)
    f_tf.paragraphs[0].font.color.rgb = EARTH

    # 2) El problema
    add_slide_title_bullets(prs, "El problema", [
        "Desorden de documentos y versiones",
        "Cálculos dispersos (Excel, correos, notas)",
        "Tiempo perdido buscando información y rehaciendo tareas",
    ])

    # 3) La propuesta de valor
    add_slide_title_bullets(prs, "La propuesta de valor", [
        "Un asistente que centraliza todo por propiedad",
        "Entiende lenguaje natural (texto/voz)",
        "Estructura el trabajo en: Documentos, Números y Resumen",
    ])

    # 4) Cómo se usa (30s)
    add_slide_title_bullets(prs, "Cómo se usa (30s)", [
        "Hablas con el asistente (chat o voz)",
        "Subes un documento; el asistente lo coloca en su sitio",
        "Rellenas números; el asistente calcula y exporta a Excel",
        "Pides un ‘resumen de la propiedad’ y lo genera en PDF/PPT",
    ])

    # 5) Esquema por propiedad
    ascii_prop = (
        "[Propiedad: Casa Demo]\n"
        "  ├─ Documentos (carpetas y subcarpetas)\n"
        "  │    • Compra: Escritura notarial, Registro público, Impuestos…\n"
        "  │    • Reforma: Licencia de obra, Contrato arquitecto, Mapas de nivel…\n"
        "  ├─ Números (tabla de variables)\n"
        "  │    • Precio de venta, Costes, Impuestos, Honorarios…\n"
        "  │    • Cálculo de totales y escenarios\n"
        "  └─ Resumen (PDF/PPT)\n"
        "       • Índice, Fotos demo, Executive summary, Mapa,\n"
        "         Tabla de números, Gráfico en cascada, Fechas clave"
    )
    add_slide_title_text(prs, "Esquema por propiedad", ascii_prop, mono=True)

    # 6) Flujo de trabajo
    flow = (
        "Crear/Seleccionar Propiedad\n"
        "        ↓\n"
        "Subir Documentos → Consultar/Resumir (Q&A con citas)\n"
        "        ↓\n"
        "Rellenar Números → Calcular → Excel\n"
        "        ↓\n"
        "Generar Resumen (PDF/PPT) → Compartir por email"
    )
    add_slide_title_text(prs, "Flujo de trabajo", flow, mono=True)

    # 7) Documentos
    add_slide_title_bullets(prs, "Documentos (qué aporta)", [
        "Estructura clara por carpetas: encuentras todo en segundos",
        "Subida guiada: propone ubicación y pide confirmación",
        "Preguntas sobre documentos: responde con citas (de dónde salió)",
    ])

    # 8) Números
    add_slide_title_bullets(prs, "Números (qué aporta)", [
        "Una tabla única para todos los valores",
        "Te dice qué falta; rellenas por voz/texto y calcula",
        "Exporta a Excel para compartir o presentar",
    ])

    # 9) Resumen de la propiedad
    add_slide_title_bullets(prs, "Resumen de la propiedad (PDF/PPT)", [
        "Documento listo para compartir",
        "Se alimenta de Documentos + Números (sin inventar)",
        "Estructura fija: índice, fotos demo, mapa, números y gráfico",
    ])

    # 10) Conversación natural
    add_slide_title_bullets(prs, "Conversación natural (texto/voz)", [
        "‘¿Qué documentos faltan?’ ‘Sube este contrato aquí’",
        "‘Pon precio de venta a 350.000€ y calcula’",
        "‘Genera la ficha resumen y envíamela por email’",
    ])

    # 11) Ejemplo rápido
    add_slide_title_bullets(prs, "Ejemplo rápido de uso (storyboard)", [
        "‘Crea Casa Demo 6 en Calle Alameda 22’",
        "‘Subo escritura notarial’",
        "‘¿Qué números me faltan?’ → ‘Pon impuestos 7%’",
        "‘Calcula’ → ‘Genera resumen en PDF’ → ‘Envíalo por email’",
    ])

    # 12) Beneficios
    add_slide_title_bullets(prs, "Beneficios para el cliente", [
        "Orden y trazabilidad: todo por propiedad",
        "Ahorro de tiempo (subida guiada, cálculos, resúmenes)",
        "Mejor calidad: respuestas con evidencia y resultados consistentes",
    ])

    # 13) Seguridad y control
    add_slide_title_bullets(prs, "Seguridad y control", [
        "Acceso seguro a datos y archivos",
        "Eliminación ‘suave’ (se oculta, no se pierde por error)",
        "Revisión y confirmación antes de subir/enviar",
    ])

    # 14) Demo en 5 minutos
    add_slide_title_bullets(prs, "Demo en 5 minutos (guion)", [
        "Crear/seleccionar propiedad",
        "Subir 1 documento (con propuesta de ubicación)",
        "Mostrar ‘documentos subidos’",
        "Rellenar 1–2 números y ‘calcular’",
        "Generar ‘resumen de la propiedad’ y mostrar link",
        "Opcional: ‘enviarlo por email’",
    ])

    # 15) Próximos pasos
    add_slide_title_bullets(prs, "Próximos pasos", [
        "Pilotar con 1–2 propiedades reales",
        "Ajustar plantillas y vistas necesarias",
        "Plan de adopción (equipo, tiempos, hitos)",
    ])

    return prs


def main():
    prs = build_client_deck()
    out_path = "client_deck_rama.pptx"
    prs.save(out_path)
    print(f"✅ Deck creado: {out_path}")


if __name__ == "__main__":
    main()


