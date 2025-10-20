from __future__ import annotations
"""
Generate only 3 improved slides for the client deck:
  1) Esquema por propiedad (visual)
  2) Plantillas: Documentos (estructura) + Números (clave)
  3) Flujo de trabajo (línea de tiempo)

Output: client_deck_rama_3.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from collections import OrderedDict
try:
    # Prefer real list from codebase (full, exact names)
    from tools.docs_tools import DOC_GROUPS, KEYWORD_TO_DOCNAME  # type: ignore
except Exception:
    # Fallback to static map if env not configured
    DOC_GROUPS = {
        "Compra": [
            "escritura notarial", "escritura", "registro publico", "registro",
            "arras", "impuestos", "impuesto", "contrato privado", "itp", "iba",
            "comentario sobre impuestos",
        ],
        "Reforma:Docs diseño": [
            "contrato arquitecto", "contrato aparejador", "mapas de nivel", "mapas",
            "planos del terreno", "planos arquitecto", "planos de la casa", "planos",
            "licencia obra", "licencia", "arquitecto", "aparejador",
        ],
        "Reforma:Docs obra": ["contrato constructor", "constructor"],
        "Reforma:Docs facturas": [
            "factura fontaneria", "factura electricista", "factura calefaccion",
            "factura carpinteria", "factura diseño", "factura",
        ],
        "Reforma:Docs registro obra nueva": ["registro documento", "documento de impuestos"],
        "Venta": ["certificacion"],
    }
    KEYWORD_TO_DOCNAME = {
        "escritura notarial": "Escritura notarial",
        "escritura": "Escritura notarial",
        "registro publico": "Registro publico",
        "registro": "Registro publico",
        "arras": "Arras",
        "impuestos": "Impuestos",
        "impuesto": "Impuestos",
        "contrato privado": "Contrato privado",
        "comentario sobre impuestos": "Comentario sobre impuestos ITP/IBA",
        "itp": "Comentario sobre impuestos ITP/IBA",
        "iba": "Comentario sobre impuestos ITP/IBA",
        "contrato arquitecto": "Contrato arquitecto",
        "contrato aparejador": "Contrato aparejador",
        "mapas de nivel": "Mapas de nivel",
        "mapas": "Mapas de nivel",
        "planos del terreno": "Planos del terreno",
        "planos arquitecto": "Planos arquitecto/de la casa",
        "planos de la casa": "Planos arquitecto/de la casa",
        "planos": "Planos arquitecto/de la casa",
        "licencia obra": "Licencia obra",
        "licencia": "Licencia obra",
        "contrato constructor": "Contrato constructor",
        "constructor": "Contrato constructor",
        "factura fontaneria": "Factura fontaneria",
        "factura electricista": "Factura electricista",
        "factura calefaccion": "Factura calefaccion",
        "factura carpinteria": "Factura carpinteria",
        "factura diseño": "Factura diseño",
        "factura": "Factura diseño",
        "registro documento": "Registro documento",
        "documento de impuestos": "Documento de impuestos",
        "certificacion": "Certificacion",
    }


GREEN = RGBColor(0x3d, 0x74, 0x35)
EARTH = RGBColor(0xc5, 0xac, 0x85)
SAND = RGBColor(0xf0, 0xec, 0xe2)
INK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide, color: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(11.4), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = "Helvetica Neue"


def card(slide, x, y, w, h):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = EARTH
    return rect


def slide_esquema_propiedad(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(s, SAND)
    title(s, "Esquema por propiedad")

    # Header chip "Propiedad"
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(2.6), Inches(0.6))
    chip.fill.solid(); chip.fill.fore_color.rgb = GREEN; chip.line.fill.background()
    ct = chip.text_frame; ct.text = "Propiedad"; ct.paragraphs[0].font.color.rgb = WHITE; ct.paragraphs[0].font.size = Pt(16); ct.paragraphs[0].alignment = PP_ALIGN.CENTER; ct.paragraphs[0].font.name = "Helvetica Neue"

    # Three cards side by side
    docs = card(s, 0.9, 2.4, 3.7, 4.3)
    nums = card(s, 4.85, 2.4, 3.7, 4.3)
    resume = card(s, 8.8, 2.4, 3.5, 4.3)

    # Documentos card
    t1 = docs.text_frame; t1.clear()
    h = t1.paragraphs[0]; h.text = "Documentos"; h.font.size = Pt(22); h.font.bold = True; h.font.color.rgb = GREEN; h.font.name = "Helvetica Neue"
    for line in [
        "Compra",
        "Reforma / Docs diseño",
        "Reforma / Docs obra",
        "Reforma / Facturas",
        "Reforma / Registro obra nueva",
        "Venta",
    ]:
        p = t1.add_paragraph(); p.text = f"• {line}"; p.font.size = Pt(16); p.font.color.rgb = INK; p.font.name = "Helvetica Neue"

    # Números card
    t2 = nums.text_frame; t2.clear()
    h2 = t2.paragraphs[0]; h2.text = "Números"; h2.font.size = Pt(22); h2.font.bold = True; h2.font.color.rgb = GREEN; h2.font.name = "Helvetica Neue"
    for line in [
        "• Precio de venta",
        "• Costes construcción",
        "• Impuestos",
        "• Honorarios",
        "• Escenarios y cálculo",
    ]:
        p = t2.add_paragraph(); p.text = line; p.font.size = Pt(16); p.font.color.rgb = INK; p.font.name = "Helvetica Neue"

    # Resumen card
    t3 = resume.text_frame; t3.clear()
    h3 = t3.paragraphs[0]; h3.text = "Resumen"; h3.font.size = Pt(22); h3.font.bold = True; h3.font.color.rgb = GREEN; h3.font.name = "Helvetica Neue"
    for line in [
        "• Índice",
        "• Fotos demo",
        "• Executive summary",
        "• Mapa",
        "• Tabla de números",
        "• Gráfico cascada",
        "• Fechas clave",
    ]:
        p = t3.add_paragraph(); p.text = line; p.font.size = Pt(16); p.font.color.rgb = INK; p.font.name = "Helvetica Neue"
    return s


def _canon_names(keywords: list[str]) -> list[str]:
    seen = OrderedDict()
    for kw in keywords:
        name = KEYWORD_TO_DOCNAME.get(kw, kw.title())
        seen[name] = True
    return list(seen.keys())


def slide_plantillas(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(s, SAND)
    title(s, "Plantillas: Documentos y Números")

    # Left column: full document groups (Compra + Reforma diseño/obra)
    c1 = card(s, 0.9, 1.8, 5.4, 4.9)
    t1 = c1.text_frame; t1.clear()
    h1 = t1.paragraphs[0]; h1.text = "Documentos"; h1.font.size = Pt(24); h1.font.bold = True; h1.font.color.rgb = GREEN; h1.font.name = "Helvetica Neue"

    left_groups = ["Compra", "Reforma:Docs diseño", "Reforma:Docs obra"]
    for g in left_groups:
        disp = g.replace(":", " — ")
        p = t1.add_paragraph(); p.text = disp; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = INK; p.font.name = "Helvetica Neue"
        names = _canon_names(DOC_GROUPS.get(g, []))
        if names:
            p2 = t1.add_paragraph(); p2.text = " • " + " · ".join(names)
            p2.font.size = Pt(14); p2.font.color.rgb = INK; p2.font.name = "Helvetica Neue"

    # Right column: remaining document groups + Numbers
    c2 = card(s, 6.0, 1.8, 6.3, 4.9)
    t2 = c2.text_frame; t2.clear()
    h2 = t2.paragraphs[0]; h2.text = "Números"; h2.font.size = Pt(24); h2.font.bold = True; h2.font.color.rgb = GREEN; h2.font.name = "Helvetica Neue"
    for line in [
        "Variables clave: Precio de venta · Costes construcción · Impuestos · Honorarios",
        "Qué falta: lista de pendientes para completar",
        "Cálculo de totales, escenarios y Excel",
    ]:
        p = t2.add_paragraph(); p.text = "• " + line; p.font.size = Pt(16); p.font.color.rgb = INK; p.font.name = "Helvetica Neue"

    # Add remaining document groups under a subtitle
    sub = t2.add_paragraph(); sub.text = "Documentos (resto de carpetas)"; sub.font.size = Pt(18); sub.font.bold = True; sub.font.color.rgb = INK; sub.font.name = "Helvetica Neue"
    right_groups = ["Reforma:Docs facturas", "Reforma:Docs registro obra nueva", "Venta"]
    for g in right_groups:
        disp = g.replace(":", " — ")
        p = t2.add_paragraph(); p.text = disp; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = INK; p.font.name = "Helvetica Neue"
        names = _canon_names(DOC_GROUPS.get(g, []))
        if names:
            p2 = t2.add_paragraph(); p2.text = " • " + " · ".join(names)
            p2.font.size = Pt(13); p2.font.color.rgb = INK; p2.font.name = "Helvetica Neue"
    return s


def slide_flujo(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(s, SAND)
    title(s, "Flujo de trabajo")

    # Build a horizontal timeline of chevrons
    steps = [
        "Crear/Seleccionar",
        "Subir documentos",
        "Q&A documentos",
        "Números · Calcular · Excel",
        "Resumen PDF/PPT",
        "Compartir por email",
    ]
    x = 0.9
    for i, label in enumerate(steps):
        w = 1.9 if i not in (3,) else 2.6
        shape = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(2.4), Inches(w), Inches(1.2))
        shape.fill.solid(); shape.fill.fore_color.rgb = GREEN if i % 2 == 0 else EARTH
        shape.line.fill.background()
        tf = shape.text_frame; tf.clear(); p = tf.paragraphs[0]
        p.text = label; p.font.size = Pt(14); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Helvetica Neue"
        x += (w - 0.2)
    return s


def main():
    prs = Presentation()
    slide_esquema_propiedad(prs)
    slide_plantillas(prs)
    slide_flujo(prs)
    out = "client_deck_rama_3_keynote.pptx"
    prs.save(out)
    print(f"✅ Deck 3 slides creado: {out}")


if __name__ == "__main__":
    main()


