from __future__ import annotations
"""
Create a refined 5-slide client deck for RAMA Agentic AI.
Focus: strong visuals, balanced layout, large typography.
Output: client_deck_rama_5.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# Palette
GREEN_DARK = RGBColor(0x1f, 0x5c, 0x3a)
GREEN_MID = RGBColor(0x3d, 0x74, 0x35)
EARTH = RGBColor(0xc5, 0xac, 0x85)
SAND = RGBColor(0xf0, 0xec, 0xe2)
INK = RGBColor(0x18, 0x18, 0x18)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(slide, SAND)

    # Big centered title card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.8), Inches(11), Inches(4))
    card.fill.solid()
    card.fill.fore_color.rgb = GREEN_DARK
    card.line.fill.background()

    tf = card.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    s = tf.add_paragraph()
    s.text = subtitle
    s.font.size = Pt(20)
    s.font.color.rgb = EARTH
    s.alignment = PP_ALIGN.CENTER
    return slide


def two_column_text(slide, left_title: str, left_points: list[str], right_title: str, right_points: list[str]):
    # Left header
    lh = slide.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(5.8), Inches(0.6))
    lht = lh.text_frame; lht.clear()
    lp = lht.paragraphs[0]; lp.text = left_title; lp.font.size = Pt(28); lp.font.bold = True; lp.font.color.rgb = GREEN_MID

    # Right header
    rh = slide.shapes.add_textbox(Inches(6.6), Inches(1.2), Inches(6.0), Inches(0.6))
    rht = rh.text_frame; rht.clear()
    rp = rht.paragraphs[0]; rp.text = right_title; rp.font.size = Pt(28); rp.font.bold = True; rp.font.color.rgb = GREEN_MID

    # Left bullets
    lb = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(5.8), Inches(4.8))
    lbt = lb.text_frame; lbt.clear()
    for i, t in enumerate(left_points):
        p = lbt.add_paragraph() if i else lbt.paragraphs[0]
        p.text = t; p.font.size = Pt(22); p.font.color.rgb = INK

    # Right bullets
    rb = slide.shapes.add_textbox(Inches(6.6), Inches(1.9), Inches(6.0), Inches(4.8))
    rbt = rb.text_frame; rbt.clear()
    for i, t in enumerate(right_points):
        p = rbt.add_paragraph() if i else rbt.paragraphs[0]
        p.text = t; p.font.size = Pt(22); p.font.color.rgb = INK


def ascii_block(slide, title: str, ascii_text: str):
    # Title
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.8))
    tft = tb.text_frame; tft.clear()
    tp = tft.paragraphs[0]; tp.text = title; tp.font.size = Pt(30); tp.font.bold = True; tp.font.color.rgb = GREEN_MID

    # ASCII block in rounded rectangle
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = EARTH
    tx = box.text_frame; tx.clear()
    p = tx.paragraphs[0]
    p.text = ascii_text
    p.font.name = "Consolas"; p.font.size = Pt(16); p.font.color.rgb = INK


def build_deck() -> Presentation:
    prs = Presentation()

    # 1. Título (grande)
    title_slide(prs, "RAMA Agentic AI", "Gestión integral de propiedades con IA")

    # 2. Introducción / Problema
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(s2, SAND)
    two_column_text(
        s2,
        left_title="Introducción",
        left_points=[
            "Asistente que centraliza el trabajo por propiedad",
            "Conversación natural (texto/voz)",
            "Tres plantillas: Documentos, Números y Resumen",
        ],
        right_title="Problema",
        right_points=[
            "Documentos y versiones dispersas",
            "Cálculos repartidos (Excel, correos, notas)",
            "Tiempo perdido y decisiones con información incompleta",
        ],
    )

    # 3. Esquema de propiedad
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(s3, SAND)
    ascii_prop = (
        "[Propiedad]\n"
        "  ├─ Documentos (carpetas y subcarpetas)\n"
        "  │    • Compra: Escritura notarial, Registro público, Impuestos…\n"
        "  │    • Reforma: Licencia obra, Contrato arquitecto, Mapas de nivel…\n"
        "  ├─ Números (tabla de variables)\n"
        "  │    • Precio de venta, Costes, Impuestos, Honorarios…\n"
        "  └─ Resumen (PDF/PPT)\n"
        "       • Índice, Fotos demo, Executive summary, Mapa,\n"
        "         Tabla de números, Gráfico en cascada, Fechas clave"
    )
    ascii_block(s3, "Esquema por propiedad", ascii_prop)

    # 4. Plantillas (Documentos + Números)
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(s4, SAND)
    two_column_text(
        s4,
        left_title="Plantilla Documentos",
        left_points=[
            "Compra: Escritura notarial, Registro público, Impuestos, Contrato privado, Comentarios ITP/IBA",
            "Reforma / Docs diseño: Contrato arquitecto, Contrato aparejador, Licencia obra, Mapas de nivel",
            "Estructura clara por carpetas y subcarpetas",
        ],
        right_title="Plantilla Números",
        right_points=[
            "Variables clave: Precio de venta, Costes construcción, Impuestos, Honorarios",
            "Qué falta: lista de pendientes para completar",
            "Cálculo de totales, escenarios y exportación a Excel",
        ],
    )

    # 5. Flujo general
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(s5, SAND)
    flow = (
        "Crear/Seleccionar Propiedad → Subir Documentos (propuesta y confirmación)\n"
        "→ Consultar/Resumir (Q&A con citas)\n"
        "→ Rellenar Números → Calcular → Excel\n"
        "→ Generar Resumen (PDF/PPT) → Compartir por email"
    )
    ascii_block(s5, "Flujo de trabajo", flow)

    return prs


def main():
    prs = build_deck()
    out = "client_deck_rama_5.pptx"
    prs.save(out)
    print(f"✅ Deck 5 slides creado: {out}")


if __name__ == "__main__":
    main()


