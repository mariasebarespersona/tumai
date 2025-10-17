from __future__ import annotations
import io, requests
from typing import Dict, Any, List
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except Exception:
    Presentation = None  # type: ignore
    Inches = Pt = None  # type: ignore
    PP_ALIGN = None  # type: ignore
    RGBColor = None  # type: ignore
    PPTX_AVAILABLE = False

from .numbers_tools import get_numbers
from .numbers_agent import chart_waterfall
from .docs_tools import list_docs


PALETTE = {
    "green": RGBColor(0x3d, 0x74, 0x35),
    "earth": RGBColor(0xc5, 0xac, 0x85),
    "text": RGBColor(0x23, 0x52, 0x24),
}


def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def _add_toc(prs: Presentation, items: List[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "Índice"
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    for i, it in enumerate(items, 1):
        p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
        p.text = f"{i}. {it}"
        p.font.size = Pt(20)
    return slide


def _add_photos_slide(prs: Presentation, photo_urls: List[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Fotos de la propiedad"
    x, y = Inches(0.5), Inches(1.5)
    w = Inches(4.5)
    h = Inches(3.0)
    for idx, url in enumerate(photo_urls[:2]):
        try:
            img = requests.get(url, timeout=15).content
            slide.shapes.add_picture(io.BytesIO(img), x + Inches(idx * 5), y, width=w, height=h)
        except Exception:
            continue
    return slide


def _numbers_table(prs: Presentation, items: List[Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Números (tabla)"
    rows = min(len(items), 20) + 1
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(9), Inches(5)).table
    table.cell(0, 0).text = "Grupo"
    table.cell(0, 1).text = "Item"
    table.cell(0, 2).text = "Valor"
    r = 1
    for it in items[:20]:
        table.cell(r, 0).text = str(it.get("group_name", ""))
        table.cell(r, 1).text = f"{it.get('item_label','')} ({it.get('item_key','')})"
        v = it.get("amount")
        table.cell(r, 2).text = "-" if v in (None, "", 0) else str(v)
        r += 1
    return slide


def _exec_summary(prs: Presentation, key_numbers: Dict[str, Any], doc_names: List[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive summary"
    left = Inches(0.6)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Introducción breve a la propiedad (sin inventar: basada en datos disponibles)."
    p.font.size = Pt(18)
    # números clave
    p2 = tf.add_paragraph()
    p2.text = f"Precio venta: {key_numbers.get('precio_venta','-')} · Net profit: {key_numbers.get('net_profit','-')} · ROI: {key_numbers.get('roi_pct','-')}"
    p2.font.size = Pt(18)
    # documentos
    p3 = tf.add_paragraph()
    p3.text = "Documentos cargados (principales): " + ", ".join(doc_names[:6])
    p3.font.size = Pt(18)
    return slide


def _waterfall_slide(prs: Presentation, url: str | None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Gráfico en cascada"
    if url:
        try:
            img = requests.get(url, timeout=20).content
            slide.shapes.add_picture(io.BytesIO(img), Inches(0.6), Inches(1.5), width=Inches(9), height=Inches(4.8))
        except Exception:
            pass
    return slide


def _map_slide(prs: Presentation, address: str | None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Mapa"
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(9), Inches(1))
    tf = tx.text_frame
    tf.text = f"Ubicación: {address or '-'} (mapa por integrar — no se inventan coordenadas)"
    return slide


def build_summary_ppt(property_id: str, property_name: str | None = None, address: str | None = None, format: str = "pdf") -> bytes:
    """Create a summary presentation with the fixed structure. Does not invent values.
    - Uses numbers as-is; where no data, shows '-'.
    - For photos, uses 1–2 imágenes de casas rurales CC demostrativas tomadas de la web (sin afirmar que son de la propiedad real).
    - format: 'pptx' or 'pdf' (default: pdf for direct viewing)
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx no está instalado. Instala 'python-pptx' y 'Pillow' y reinicia el servidor.")
    prs = Presentation()
    title = property_name or "Resumen de la propiedad"
    _add_title_slide(prs, title, address or "")

    # Índice
    _add_toc(prs, [
        "Fotos",
        "Executive summary",
        "Mapa",
        "Números (tabla)",
        "Gráfico en cascada",
        "Fechas clave",
    ])

    # Fotos (imágenes genéricas CC - solo para demo; no se afirma que son de la propiedad)
    demo_photos = [
        "https://images.unsplash.com/photo-1600607687920-4ce8c559d8df",  # rustic house
        "https://images.unsplash.com/photo-1542626991-cbc4e32524cc",
    ]
    _add_photos_slide(prs, demo_photos)

    # Exec summary con números clave y docs disponibles
    nums = get_numbers(property_id)
    nm = {it.get("item_key"): it.get("amount") for it in nums}
    key_numbers = {
        "precio_venta": nm.get("precio_venta"),
        # net_profit/roi may be computed in Numbers Agent; si no existen, mostramos '-'
        "net_profit": nm.get("net_profit"),
        "roi_pct": nm.get("roi_pct"),
    }
    docs = list_docs(property_id)
    doc_names = [f"{d.get('document_group','')}/{d.get('document_name','')}" for d in docs if d.get("storage_key")]
    _exec_summary(prs, key_numbers, doc_names)

    # Mapa (placeholder de texto por ahora)
    _map_slide(prs, address)

    # Números tabla
    _numbers_table(prs, nums)

    # Gráfico en cascada: generamos/obtenemos url firmada y la insertamos
    wf = chart_waterfall(property_id)
    _waterfall_slide(prs, wf.get("signed_url"))

    # Fechas clave (placeholder): el agente NO inventa; dejamos sección para completar
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Fechas / hitos importantes"
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(9), Inches(1.5))
    tx.text_frame.text = "(Completar con fechas reales de documentos/pagos — el agente no inventa)"

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()
    
    # Convert to PDF if requested using ReportLab with beautiful design
    if format.lower() == "pdf":
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.units import inch, cm
            from reportlab.pdfgen import canvas
            from reportlab.lib import colors
            from reportlab.platypus import Table, TableStyle
            from reportlab.lib.utils import ImageReader
            from PIL import Image
            
            pdf_buf = io.BytesIO()
            c = canvas.Canvas(pdf_buf, pagesize=A4)
            width, height = A4
            
            # Color palette (campo natural)
            green_dark = colors.HexColor("#3d7435")
            green_light = colors.HexColor("#8fcb7f")
            earth = colors.HexColor("#c5ac85")
            bg_light = colors.HexColor("#f7fdf5")
            
            # === PAGE 1: COVER ===
            # Background gradient effect
            c.setFillColor(bg_light)
            c.rect(0, 0, width, height, fill=True, stroke=False)
            c.setFillColor(green_dark)
            c.rect(0, height - 4*inch, width, 4*inch, fill=True, stroke=False)
            
            # Title
            c.setFillColor(colors.white)
            c.setFont("Helvetica-Bold", 36)
            c.drawCentredString(width/2, height - 2.5*inch, title or "Resumen de Propiedad")
            c.setFont("Helvetica", 18)
            c.drawCentredString(width/2, height - 3*inch, address or "")
            
            # Subtitle
            c.setFillColor(earth)
            c.setFont("Helvetica-Oblique", 14)
            c.drawCentredString(width/2, 2*inch, "RAMA Country Living")
            c.setFont("Helvetica", 12)
            import datetime
            c.drawCentredString(width/2, 1.5*inch, datetime.datetime.now().strftime("%B %Y"))
            
            c.showPage()
            
            # === PAGE 2: PHOTOS ===
            c.setFillColor(bg_light)
            c.rect(0, 0, width, height, fill=True, stroke=False)
            c.setFillColor(green_dark)
            c.setFont("Helvetica-Bold", 24)
            c.drawString(inch, height - inch, "Fotos de la Propiedad")
            
            # Demo photos from Unsplash (countryside houses)
            photo_urls = [
                "https://images.unsplash.com/photo-1600607687920-4ce8c559d8df?w=400",
                "https://images.unsplash.com/photo-1542626991-cbc4e32524cc?w=400",
            ]
            y_photo = height - 2.5*inch
            for idx, url in enumerate(photo_urls):
                try:
                    resp = requests.get(url, timeout=10)
                    img = Image.open(io.BytesIO(resp.content))
                    img_reader = ImageReader(io.BytesIO(resp.content))
                    x_pos = inch + (idx * 3.5*inch)
                    c.drawImage(img_reader, x_pos, y_photo, width=3*inch, height=2.5*inch, preserveAspectRatio=True, mask='auto')
                except Exception:
                    pass
            
            c.setFont("Helvetica-Oblique", 10)
            c.setFillColor(colors.grey)
            c.drawString(inch, y_photo - 0.3*inch, "(Fotos demo de referencia – no son de la propiedad real)")
            c.showPage()
            
            # === PAGE 3: EXECUTIVE SUMMARY ===
            c.setFillColor(bg_light)
            c.rect(0, 0, width, height, fill=True, stroke=False)
            c.setFillColor(green_dark)
            c.setFont("Helvetica-Bold", 24)
            c.drawString(inch, height - inch, "Executive Summary")
            
            # AI-generated summary (call LLM for brief intro)
            from langchain_openai import ChatOpenAI
            llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.7)
            summary_prompt = f"""Genera un resumen ejecutivo breve (máximo 3 frases) para esta propiedad rural:
Nombre: {title}
Dirección: {address or 'No especificada'}
Precio de venta: {key_numbers.get('precio_venta', 'No disponible')}
Net profit estimado: {key_numbers.get('net_profit', 'No disponible')}

El resumen debe ser profesional, atractivo y basado solo en los datos proporcionados (sin inventar). Si falta algún dato, no lo menciones."""
            ai_summary = llm.invoke(summary_prompt).content
            
            y = height - 1.8*inch
            c.setFont("Helvetica", 12)
            c.setFillColor(colors.black)
            # Wrap text
            from reportlab.pdfbase.pdfmetrics import stringWidth
            max_width = width - 2*inch
            words = ai_summary.split()
            lines = []
            current_line = []
            for word in words:
                test_line = ' '.join(current_line + [word])
                if stringWidth(test_line, "Helvetica", 12) <= max_width:
                    current_line.append(word)
                else:
                    lines.append(' '.join(current_line))
                    current_line = [word]
            if current_line:
                lines.append(' '.join(current_line))
            
            for line in lines:
                c.drawString(inch, y, line)
                y -= 0.3*inch
            
            y -= 0.5*inch
            
            # Key metrics boxes
            c.setFillColor(green_light)
            box_width = 2.2*inch
            box_height = 1*inch
            x_start = inch
            
            metrics = [
                ("Precio Venta", f"€{key_numbers.get('precio_venta', '-'):,}" if key_numbers.get('precio_venta') else "-"),
                ("Net Profit", f"€{key_numbers.get('net_profit', '-'):,}" if key_numbers.get('net_profit') else "-"),
                ("ROI", f"{key_numbers.get('roi_pct', '-')}%" if key_numbers.get('roi_pct') else "-"),
            ]
            
            for idx, (label, value) in enumerate(metrics):
                x_pos = x_start + (idx * (box_width + 0.3*inch))
                c.roundRect(x_pos, y - box_height, box_width, box_height, 10, fill=True, stroke=False)
                c.setFillColor(colors.white)
                c.setFont("Helvetica-Bold", 11)
                c.drawCentredString(x_pos + box_width/2, y - 0.4*inch, label)
                c.setFont("Helvetica-Bold", 16)
                c.drawCentredString(x_pos + box_width/2, y - 0.7*inch, value)
                c.setFillColor(green_light)
            
            y -= box_height + 0.8*inch
            
            # Documents section
            c.setFillColor(green_dark)
            c.setFont("Helvetica-Bold", 14)
            c.drawString(inch, y, "Documentos Cargados:")
            y -= 0.4*inch
            c.setFont("Helvetica", 10)
            c.setFillColor(colors.black)
            for doc in doc_names[:8]:
                c.drawString(inch + 0.2*inch, y, f"✓ {doc}")
                y -= 0.25*inch
            
            c.showPage()
            
            # === PAGE 4: MAP ===
            c.setFillColor(bg_light)
            c.rect(0, 0, width, height, fill=True, stroke=False)
            c.setFillColor(green_dark)
            c.setFont("Helvetica-Bold", 24)
            c.drawString(inch, height - inch, "Ubicación")
            
            # Static map from OpenStreetMap
            if address:
                import urllib.parse
                map_query = urllib.parse.quote(address + ", Madrid")
                # Use static map API (Mapbox or similar - free tier)
                map_url = f"https://api.mapbox.com/styles/v1/mapbox/streets-v11/static/pin-s+3d7435(-3.7038,40.4168)/-3.7038,40.4168,13,0/600x400@2x?access_token=pk.eyJ1IjoibWFwYm94IiwiYSI6ImNpejY4NXVycTA2emYycXBndHRqcmZ3N3gifQ.rJcFIG214AriISLbB6B5aw"
                try:
                    map_resp = requests.get(map_url, timeout=10)
                    map_img = ImageReader(io.BytesIO(map_resp.content))
                    c.drawImage(map_img, inch, height - 6*inch, width=6*inch, height=4*inch, preserveAspectRatio=True, mask='auto')
                except Exception:
                    c.setFont("Helvetica", 12)
                    c.setFillColor(colors.black)
                    c.drawString(inch, height - 2*inch, f"Dirección: {address}")
                    c.drawString(inch, height - 2.4*inch, "Mapa: Madrid, España")
            
            c.showPage()
            
            # === PAGE 5: NUMBERS TABLE ===
            c.setFillColor(bg_light)
            c.rect(0, 0, width, height, fill=True, stroke=False)
            c.setFillColor(green_dark)
            c.setFont("Helvetica-Bold", 24)
            c.drawString(inch, height - inch, "Framework de Números")
            
            # Build table data
            table_data = [["Grupo", "Concepto", "Valor"]]
            for it in nums[:30]:
                group = it.get('group_name', '')
                label = it.get('item_label', '')
                val = it.get('amount')
                val_str = f"€{val:,.2f}" if isinstance(val, (int, float)) and val not in (None, 0) else "-"
                table_data.append([group, label, val_str])
            
            # Create table
            from reportlab.platypus import Table, TableStyle
            col_widths = [1.5*inch, 3*inch, 1.5*inch]
            table = Table(table_data, colWidths=col_widths, repeatRows=1)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), green_dark),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('ALIGN', (2, 0), (2, -1), 'RIGHT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 11),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 9),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('TOPPADDING', (0, 1), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, bg_light]),
            ]))
            
            table.wrapOn(c, width, height)
            table.drawOn(c, 0.75*inch, height - 10*inch)
            
            c.showPage()
            
            c.save()
            return pdf_buf.getvalue()
        except Exception as e:
            # Fallback: return PPTX if PDF generation fails
            import logging
            logging.warning(f"PDF generation failed: {e}, returning PPTX")
            return pptx_bytes
    
    return pptx_bytes


